"""Build GlobalMinds_Pitch_Deck.pptx — 23 slides + speaker notes."""
from __future__ import annotations
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ─── Brand ────────────────────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0B, 0x14, 0x26)
BG_PANEL = RGBColor(0x15, 0x23, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIME = RGBColor(0xC5, 0xFF, 0x4A)
GRAY = RGBColor(0x94, 0xA3, 0xB8)
GRAY_DIM = RGBColor(0x64, 0x74, 0x8B)
RED = RGBColor(0xF8, 0x71, 0x71)
GREEN = RGBColor(0x34, 0xD3, 0x99)
BLUE = RGBColor(0x60, 0xA5, 0xFA)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
PINK = RGBColor(0xF4, 0x72, 0xB6)

ROOT = Path("presentation/globalminds_microsoft_pitch_may2026")
ASSETS = ROOT / "assets"
OUT = ROOT / "GlobalMinds_Pitch_Deck.pptx"

# 16:9 widescreen
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ─── Helpers ──────────────────────────────────────────────────────────────────
def set_bg(slide, color: RGBColor = BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=18, color=WHITE,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font="Inter"):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    return tx


def add_rect(slide, left, top, width, height, fill: RGBColor = BG_PANEL,
             line: RGBColor | None = None, line_w: float = 1.0,
             rounded: bool = True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if rounded:
        # Use a small adjustment value (smaller corners) by manipulating the avLst
        try:
            sp_pr = s._element.spPr
            prst_geom = sp_pr.find(qn("a:prstGeom"))
            if prst_geom is not None:
                av_lst = prst_geom.find(qn("a:avLst"))
                if av_lst is None:
                    av_lst = etree.SubElement(prst_geom, qn("a:avLst"))
                gd = etree.SubElement(av_lst, qn("a:gd"))
                gd.set("name", "adj")
                gd.set("fmla", "val 8000")  # rounded but not too much
        except Exception:
            pass
    return s


def add_footer(slide, page: int, total: int = 23):
    add_text(slide, Inches(0.4), Inches(7.15), Inches(8), Inches(0.3),
             "Confidencial · GlobalMinds · Mayo 2026",
             size=9, color=GRAY_DIM)
    add_text(slide, Inches(11.5), Inches(7.15), Inches(1.5), Inches(0.3),
             f"{page} / {total}", size=9, color=GRAY_DIM, align=PP_ALIGN.RIGHT)


def add_lime_underline(slide, left, top, width, line_w=2.5):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(line_w))
    s.fill.solid()
    s.fill.fore_color.rgb = LIME
    s.line.fill.background()
    return s


def add_image(slide, path, left, top, width=None, height=None):
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_title_block(slide, title: str, subtitle: str | None = None):
    add_text(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7),
             title, size=30, color=WHITE, bold=True)
    add_lime_underline(slide, Inches(0.5), Inches(1.05), Inches(0.8))
    if subtitle:
        add_text(slide, Inches(0.5), Inches(1.15), Inches(12.3), Inches(0.5),
                 subtitle, size=14, color=GRAY, italic=True)


def set_notes(slide, text: str):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


# ─── Build ────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


# ── Slide 1 — Portada (full hero) ────────────────────────────────────────────
def slide_01():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_image(s, ASSETS / "diagrams/01_hero.png", Inches(0), Inches(0),
              width=SLIDE_W, height=SLIDE_H)
    set_notes(s, """[30s — Apertura]

Saludo: 'Buenos días. Soy Pedro. Con Freddy presentamos GlobalMinds — la plataforma de IA aplicada al mercado cambiario de Latinoamérica.'

Hook: 'Tenemos 18 meses construyendo, un sistema en producción HOY, y 4 oportunidades comerciales validadas con datos públicos. En 25 minutos les muestro por qué este es el momento de entrar.'

Tono: confiado, directo. No sonreír de más. Pausar 2 segundos antes de pasar.""")

# ── Slide 2 — El problema ────────────────────────────────────────────────────
def slide_02():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "El problema", "El mercado cambiario LATAM mueve billones diarios — sin IA seria")

    # 4 stat cards (2x2 grid)
    cards = [
        ("USD 13,098 M", "remesas Colombia 2025 (récord histórico) convertidas a TRM del día, sin optimización de timing",
         GREEN),
        ("17%", "de empresas no financieras colombianas usa derivados de cobertura FX (último estudio público disponible Banco de la República)",
         BLUE),
        ("USD 50–200k", "diarios mueve cada merchant top de Binance P2P Colombia — con Excel y 4 personas, cero ML",
         AMBER),
        ("Decenas de M COP", "le cuesta una operación mal cronometrada a una PYME exportadora promedio",
         RED),
    ]
    positions = [(0.5, 1.9), (6.95, 1.9), (0.5, 4.65), (6.95, 4.65)]
    for (big, desc, color), (cx, cy) in zip(cards, positions):
        add_rect(s, Inches(cx), Inches(cy), Inches(5.85), Inches(2.5),
                 fill=BG_PANEL, line=color, line_w=1.5)
        add_text(s, Inches(cx + 0.35), Inches(cy + 0.25), Inches(5.5), Inches(0.9),
                 big, size=34, bold=True, color=color)
        add_text(s, Inches(cx + 0.35), Inches(cy + 1.15), Inches(5.5), Inches(1.3),
                 desc, size=13, color=WHITE)

    add_footer(s, 2)
    set_notes(s, """[60s — El problema]

Tres datos memorizar:

1. 'Colombia recibió USD 13 mil millones en remesas en 2025 — RÉCORD HISTÓRICO. Marzo solo: USD 1,225 millones en un mes. Esto es un flujo gigante que se convierte ciegamente.'

2. 'El último estudio público del Banco de la República muestra que solo el 17% de empresas no financieras usa cobertura FX en Colombia, vs más del 60% en economías avanzadas. Hay 5 mil empresas exportadoras desprotegidas.'

3. 'En Binance P2P un solo merchant top mueve hasta USD 200 mil diarios — con un Excel y 4 personas. Cero ML, cero automatización.'

Cierre: 'Tres mercados gigantes. Cero IA seria. Ahí es donde entramos.'

OJO — si te preguntan por el 17%: 'Es el dato más reciente publicado por Banrep, Borradores 1058. Reconocemos que el universo puede haber cambiado, pero la tendencia se mantiene según estudios sectoriales más recientes.'""")

# ── Slide 3 — La oportunidad ────────────────────────────────────────────────
def slide_03():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "La oportunidad", "4 mercados sin atender · 1 sola plataforma de IA · TAM USD 50–150M LATAM")

    # 4 mercados grid
    markets = [
        ("Merchants P2P", "Cientos de merchants moviendo USD 5–200k/día sin ML", GREEN),
        ("PYMES con riesgo FX", "5,000+ exportadoras/importadoras desprotegidas", BLUE),
        ("Casas de cambio", "USD 13B/año en remesas convertidas sin optimización", AMBER),
        ("Freelancers cobrando USD", "500k–1M colombianos pierden 1.5–3% en Wise/Payoneer", PINK),
    ]
    for i, (title, desc, color) in enumerate(markets):
        col = i % 2
        row = i // 2
        cx = 0.5 + col * 6.45
        cy = 1.9 + row * 1.5
        add_rect(s, Inches(cx), Inches(cy), Inches(6.35), Inches(1.3),
                 fill=BG_PANEL, line=color, line_w=1.4)
        # bullet circle
        bx = cx + 0.3
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(bx), Inches(cy + 0.45),
                                    Inches(0.4), Inches(0.4))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = color
        bullet.line.fill.background()
        add_text(s, Inches(bx + 0.55), Inches(cy + 0.2), Inches(5.5), Inches(0.5),
                 title, size=18, bold=True, color=WHITE)
        add_text(s, Inches(bx + 0.55), Inches(cy + 0.7), Inches(5.5), Inches(0.6),
                 desc, size=12, color=GRAY)

    # bottom callout
    add_rect(s, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.3),
             fill=BG_PANEL, line=LIME, line_w=2)
    add_text(s, Inches(0.7), Inches(5.45), Inches(12), Inches(0.5),
             "TAM combinado en Latinoamérica", size=13, color=LIME, bold=True)
    add_text(s, Inches(0.7), Inches(5.85), Inches(12), Inches(0.7),
             "USD 50–150M anuales — y Microsoft / AWS buscan casos LATAM emblemáticos en fintech ML",
             size=18, color=WHITE, bold=True)

    add_footer(s, 3)
    set_notes(s, """[45s — La oportunidad]

'Hay 4 mercados gigantes en Colombia y LATAM que comparten algo: deciden con USD/COP, USD/MXN, USD/BRL — sin ML serio. Son: merchants P2P, PYMES con riesgo FX, casas de cambio que mueven remesas, y freelancers que cobran en USD.'

'Combinados representan un mercado direccionable de 50 a 150 millones de dólares anuales en LATAM. Y la cereza: Microsoft y AWS están buscando activamente casos LATAM emblemáticos en fintech con ML — somos el caso que necesitan.'

Pausa breve. Mirar a la audiencia.""")

# ── Slide 4 — ¿Qué construimos? ─────────────────────────────────────────────
def slide_04():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "¿Qué construimos?", "Sistema end-to-end de ML + RL · 18 meses · en producción hoy")

    # left big card — el motor
    add_rect(s, Inches(0.5), Inches(1.9), Inches(6), Inches(4.9),
             fill=BG_PANEL, line=LIME, line_w=2)
    add_text(s, Inches(0.85), Inches(2.15), Inches(5.5), Inches(0.6),
             "Núcleo USD/COP", size=22, bold=True, color=LIME)
    add_text(s, Inches(0.85), Inches(2.7), Inches(5.5), Inches(0.6),
             "Validado en producción", size=12, color=GRAY, italic=True)

    # bullets
    bullets = [
        "Backtest 2025 OOS: +25.63% / Sharpe 3.35 / p = 0.006",
        "82.4% win rate · 34 trades · maxDD 6.12%",
        "Walk-forward weekly retraining (anti-leakage)",
        "Regime gate (Hurst): bloqueó 13 de 14 semanas mean-reverting Q1 2026",
        "Buy & hold mismo período: −14.48%",
    ]
    by = 3.4
    for b in bullets:
        # check icon
        chk = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), Inches(by + 0.07),
                                 Inches(0.18), Inches(0.18))
        chk.fill.solid(); chk.fill.fore_color.rgb = LIME
        chk.line.fill.background()
        add_text(s, Inches(1.15), Inches(by), Inches(5.2), Inches(0.5),
                 b, size=12, color=WHITE)
        by += 0.6

    # right card — expansion
    add_rect(s, Inches(6.85), Inches(1.9), Inches(6), Inches(4.9),
             fill=BG_PANEL, line=BLUE, line_w=1.5)
    add_text(s, Inches(7.2), Inches(2.15), Inches(5.5), Inches(0.6),
             "Expansión LATAM", size=22, bold=True, color=BLUE)
    add_text(s, Inches(7.2), Inches(2.7), Inches(5.5), Inches(0.6),
             "Mismo motor — pares vecinos", size=12, color=GRAY, italic=True)

    pairs = [
        ("USD/MXN", "México · 95K bars históricos listos"),
        ("USD/BRL", "Brasil · 90K bars históricos listos"),
        ("USD/CLP", "Chile · roadmap Q3 2026"),
        ("USD/PEN", "Perú · roadmap Q4 2026"),
    ]
    py = 3.4
    for pair, desc in pairs:
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.2), Inches(py + 0.07),
                                  Inches(0.18), Inches(0.18))
        circ.fill.solid(); circ.fill.fore_color.rgb = BLUE
        circ.line.fill.background()
        add_text(s, Inches(7.5), Inches(py), Inches(2), Inches(0.5),
                 pair, size=14, bold=True, color=WHITE)
        add_text(s, Inches(9.5), Inches(py), Inches(3.3), Inches(0.5),
                 desc, size=11, color=GRAY)
        py += 0.6

    add_footer(s, 4)
    set_notes(s, """[60s — Producto]

'Lo primero que necesitan saber: esto NO es una promesa, es producto.'

Izquierda — '18 meses de desarrollo. El núcleo USD/COP está en producción hoy. Backtest 2025 — out of sample, walk-forward, auditado: 25.63% de retorno, Sharpe 3.35, p-value 0.006. 82% win rate. Y un componente que es el verdadero MVP: el regime gate que detecta cuándo NO operar — bloqueó 13 de 14 semanas mean-reverting en Q1 2026 y por eso no perdimos plata.'

Derecha — 'El mismo motor se aplica a USD/MXN, USD/BRL, USD/CLP. Tenemos 95 mil barras históricas de México y 90 mil de Brasil ya cargadas. La expansión es replicación, no R&D nuevo.'

Cierre: 'Mismo CAPEX, 4× mercado.'""")

# ── Slide 5 — Arquitectura ──────────────────────────────────────────────────
def slide_05():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Arquitectura: 1 motor → 4 productos",
                    "Enterprise-grade · cloud-ready (Azure / AWS / GCP)")
    # full-width image
    add_image(s, ASSETS / "diagrams/05_architecture.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 5)
    set_notes(s, """[45s — Arquitectura]

Recorrer de arriba a abajo, con calma:

'Arriba: ingestamos 4 fuentes — mercado FX en 5-min, 40 indicadores macro, 5 fuentes de noticias bilingües, y análisis IA con GPT-4o + Claude.'

'En el centro está el motor: ensemble Ridge + BayesianRidge + XGBoost, con un regime gate basado en Hurst que decide cuándo NO operar, y entrenamiento walk-forward semanal — sin leakage.'

'Abajo: el mismo motor genera señales universales que alimentan los 4 productos comerciales. Mismo código, 4 monetizaciones.'

'Stack a la derecha: Airflow, MLflow, Postgres con TimescaleDB, FastAPI, Next.js, Docker. Listo para correr en Azure ML, OpenAI, Container Apps — exactamente las piezas que les interesa promover a Microsoft.'""")

# ── Slide 6 — Demo dashboard (3 screenshots) ────────────────────────────────
def slide_06():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Producto en producción · capturas reales",
                    "Dashboard Next.js · datos reales · sin mockups")
    # 3 screenshots side-by-side
    ss_w = Inches(4.05)
    ss_y = Inches(1.85)
    ss_h = Inches(2.35)
    paths = [
        ("forecasting.png", "FORECASTING — model zoo"),
        ("dashboard.png", "BACKTEST 2025 — KPIs + trades"),
        ("production.png", "PRODUCTION — 2026 YTD live"),
    ]
    starts = [0.45, 4.65, 8.85]
    for (fn, label), x in zip(paths, starts):
        add_image(s, ASSETS / "screenshots" / fn, Inches(x), ss_y, width=ss_w)
        # caption box
        add_rect(s, Inches(x), Inches(4.35), ss_w, Inches(0.55),
                 fill=BG_PANEL, line=LIME, line_w=1)
        add_text(s, Inches(x), Inches(4.45), ss_w, Inches(0.4),
                 label, size=11, bold=True, color=LIME, align=PP_ALIGN.CENTER)

    # bottom highlights
    add_rect(s, Inches(0.45), Inches(5.15), Inches(12.45), Inches(1.7),
             fill=BG_PANEL, line=GRAY_DIM, line_w=1)
    highlights = [
        ("47", "API routes activas"),
        ("8", "páginas Next.js"),
        ("25+", "servicios Docker"),
        ("27", "DAGs Airflow productivos"),
    ]
    for i, (val, lbl) in enumerate(highlights):
        x = 0.7 + i * 3.05
        add_text(s, Inches(x), Inches(5.35), Inches(2.9), Inches(0.7),
                 val, size=32, bold=True, color=LIME, align=PP_ALIGN.CENTER)
        add_text(s, Inches(x), Inches(6.15), Inches(2.9), Inches(0.5),
                 lbl, size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_footer(s, 6)
    set_notes(s, """[40s — Demo dashboard]

'Esto es producto real. Ninguna captura es mockup.'

Señalar cada captura:

1. Izquierda — Forecasting: 'Aquí ven 9 modelos corriendo con backtests walk-forward por horizonte. Es nuestro model zoo interno.'

2. Centro — Backtest 2025: 'Esta es la pantalla que el equipo de aprobación humana revisa antes de promover un modelo a producción. Vote 2 de 2.'

3. Derecha — Production: 'Y aquí está la operación 2026 en vivo. El regime gate bloqueando trades cuando no debe operar — eso es lo que nos protegió.'

Cierre: '47 API routes, 8 páginas, 27 DAGs en Airflow. Esto es enterprise-grade, no un script de notebook.'""")

# ── Slide 7 — Storytelling de performance ───────────────────────────────────
def slide_07():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Performance auditada — backtest 2025 OOS",
                    "Out-of-sample real · walk-forward · sin look-ahead bias")

    add_image(s, ASSETS / "charts/07_equity_curve_2025.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 7)
    set_notes(s, """[60s — Performance]

'Esto es lo que separa GlobalMinds de cualquier promesa: datos verificables.'

'Línea verde: nuestro modelo. $10,000 iniciales se convierten en $12,563 en 12 meses. Línea roja punteada: buy and hold del USD/COP — termina en $8,552. Diferencia: 40 puntos porcentuales de alpha en un año.'

'Sharpe 3.35 — Sharpe institucional. P-value 0.006 — la probabilidad de que esto sea suerte es 6 en mil. Win rate 82.4% en 34 trades.'

Caso acertado: 'Semana 18 — Hurst trending fuerte, DXY arriba: SHORT entró, take profit en 4 días, +1.4%.'
Caso fallido (transparencia): 'Semana 36 — shock político no previsto, hard stop, −2.1%. Esto pasa. Lo importante es que el sizing y el regime gate evitaron que fuera 5×.'

Si te preguntan por overfitting: 'Walk-forward semanal, expanding window desde 2020. Anti-leakage estricto: macro features T-1, norm_stats solo del train split. Auditado por 10 agentes especializados.'""")

# ── Slide 8 — Las 4 oportunidades ───────────────────────────────────────────
def slide_08():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Las 4 oportunidades validadas",
                    "1 motor de IA · 4 mercados sin atender · USD 50–150M TAM LATAM")
    add_image(s, ASSETS / "diagrams/08_opportunities.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 8)
    set_notes(s, """[60s — 4 oportunidades]

'Cuatro líneas de negocio que comparten el mismo motor de IA, mismo CAPEX, 4× revenue.'

Recorrer cada card en orden:

1. Bot P2P — 'merchants en Binance moviendo cientos de miles diarios con Excel. USD 99 a 1,500/mes según volumen. Cash flow más rápido. 15 a 25k MRR a 18 meses.'

2. Hedging PYMES — 'Kantox para LATAM. Kantox vendió a Visa por 175 millones de euros en 2025. 5 mil PYMES desprotegidas. 60k MRR a 24 meses.'

3. Remesas B2B — 'Layer de IA sobre 13 mil millones de dólares anuales. No tocamos al consumidor — vendemos a casas de cambio. 10 a 25k MRR.'

4. App freelancers — 'medio millón a un millón de colombianos cobrando en USD. Wise les cobra 1.5%, nosotros les regresamos 1 a 2% con timing IA. 40k MRR a 24 meses.'

Cierre: 'No tenemos que escoger una. La arquitectura las soporta todas — el orden de salida lo dicta la velocidad de cash flow.'""")

# ── Slide 9 — Datos que validan cada mercado ────────────────────────────────
def slide_09():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Datos públicos que validan cada mercado",
                    "Cifras verificables · fuentes oficiales · ningún supuesto")

    rows = [
        ("P2P", "Binance cero fees · spread 1.3–3.9% bruto Top1 vs Top10 · merchants top USD 50–200k/día con Excel",
         GREEN),
        ("Hedging", "Solo 17% de empresas no financieras COL usa derivados FX (estudio Banco de la República). Comparable: Kantox vendida a Visa por €175M (2025).",
         BLUE),
        ("Remesas", "USD 13,098M en 2025 — récord histórico (Banco de la República). 1.6× la IED en Q1 2026. 2.1M receptores. 53% origen EE.UU.",
         AMBER),
        ("Freelancers", "500k–1M colombianos cobrando USD desde Upwork, Deel, Fiverr, Remote. Wise cobra 1.5%, Payoneer 2–3%. Bre-B (>100M llaves) reduce fricción de salida a COP.",
         PINK),
    ]
    yt = 1.85
    for tag, desc, color in rows:
        # left tag block
        add_rect(s, Inches(0.5), Inches(yt), Inches(2.5), Inches(1.15),
                 fill=color, line=None, line_w=0)
        add_text(s, Inches(0.5), Inches(yt + 0.3), Inches(2.5), Inches(0.6),
                 tag, size=22, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER)
        # description block
        add_rect(s, Inches(3.1), Inches(yt), Inches(9.8), Inches(1.15),
                 fill=BG_PANEL, line=GRAY_DIM, line_w=0.5)
        add_text(s, Inches(3.4), Inches(yt + 0.2), Inches(9.4), Inches(0.85),
                 desc, size=12, color=WHITE)
        yt += 1.3

    add_footer(s, 9)
    set_notes(s, """[45s — Datos que validan]

Pasar por cada fila rápido, dejando que las cifras hablen:

P2P: 'Binance no cobra fees. El spread bruto entre Top 1 y Top 10 es de 1.3 a 3.9%. Esa es la captura.'

Hedging: 'Banco de la República — el último estudio público disponible — muestra 17% de penetración. Y Visa pagó 175 millones de euros por Kantox el año pasado, que hace exactamente esto.'

Remesas: 'Banco de la República otra vez: 13 mil millones en 2025, récord histórico. Por primera vez en 20 años las remesas superaron a la IED. 53% viene de Estados Unidos — flujo súper estable.'

Freelancers: 'Medio millón a un millón de colombianos. Wise les cobra 1.5%. Bre-B con 100 millones de llaves activas elimina la fricción de salida. La infra ya está — nosotros solo orquestamos.'

Cada cifra tiene fuente. Slide 23 las lista todas.""")

# ── Slide 10 — Comparables ──────────────────────────────────────────────────
def slide_10():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Comparables internacionales",
                    "Validación con valuaciones reales — ninguno opera Colombia con esta propuesta")
    add_image(s, ASSETS / "diagrams/10_comparables.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 10)
    set_notes(s, """[40s — Comparables]

'No estamos inventando una categoría — la categoría existe y está validada por dinero institucional.'

Recorrer rápido:
- 'Kantox: hedging FX para PYMES. Visa la compró por 175 millones de euros en 2025.'
- 'Bound, Neo, Pangea, TreasurUp: todas con Series A activas en el mismo segmento.'
- 'Wise: 11 mil millones de dólares de valuación pública. Remitly: 3 mil millones. Deel: 12 mil millones.'
- 'XTX Markets: el referente global de ML aplicado a market making en FX.'

Cierre fuerte: 'NINGUNO de estos opera Colombia con esta propuesta combinada. Nos pueden copiar — pero nosotros llegamos primero, con el motor ya probado.'""")

# ── Slide 11 — 4 motores ingreso ────────────────────────────────────────────
def slide_11():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "4 motores de ingreso · 1 solo motor de IA",
                    "Mismo CAPEX · 4× revenue · diversificación · cross-sell")
    add_image(s, ASSETS / "diagrams/11_revenue_streams.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 11)
    set_notes(s, """[40s — Modelo monetización]

'Esto es el corazón del modelo financiero: un solo motor — el que ya construimos — alimenta 4 streams independientes.'

'Por qué es poderoso: si el bot P2P se demora en escalar, hedging tira. Si hedging requiere ciclo largo de venta corporativo, freelancers viral. Diversificación de riesgo dentro de la misma compañía.'

'Cross-sell: el merchant P2P que trabajamos hoy nos da el case study para vender hedging a su importadora hermana.'

'Margen SaaS típico 70 a 85%. Y todos los productos son BYOK — Bring Your Own Keys — el cliente paga sus propias APIs de Binance, AWS, etc. Nosotros no asumimos costo variable.'""")

# ── Slide 12 — Pricing ──────────────────────────────────────────────────────
def slide_12():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Pricing por línea de producto",
                    "Anclado en comparables internacionales · tiered por volumen")
    add_image(s, ASSETS / "diagrams/12_pricing.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 12)
    set_notes(s, """[35s — Pricing]

'Precios anclados en lo que ya cobran los comparables: Kantox arranca en €299/mes, Bound similar. Wise saca tarifa transaccional pura — nosotros agregamos tier.'

Recorrer rápido:
- 'P2P: 99 a 199 para volumen pequeño, 499 a 1,500 para Whales con performance fee.'
- 'Hedging: tier PYME, tier Mid, tier Corporate con revenue share con banco partner.'
- 'Remesas: floor de 2 mil al mes y 10 a 20% de rev share sobre el alpha capturado.'
- 'Freelancers: suscripción + transaccional, modelo Wise + Spotify.'

'Margen 70 a 85% es estándar SaaS. BYOK significa cero costo variable. Estos números no requieren capex — solo customer acquisition.'""")

# ── Slide 13 — Bot P2P punta de lanza ────────────────────────────────────────
def slide_13():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Bot P2P — Punta de lanza ⭐",
                    "El producto con cash flow más rápido y mercado más virgen")
    add_image(s, ASSETS / "diagrams/13_p2p_lead.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 13)
    set_notes(s, """[60s — Bot P2P]

'Si solo me dejaran vender una de las cuatro, sería esta. Y se las explico.'

'Mercado: cero fees de plataforma en Binance P2P. El spread bruto Top 1 vs Top 10 es de 1.3 a 3.9%. Un merchant top mueve hasta USD 200 mil por día.'

'Estado actual: estos merchants operan con Excel, 4 personas vigilando pantallas, manualmente. Un solo error humano les cuesta el spread del día.'

'Nuestro bot agrega 5 cosas: detección de chargebacks con ML, ajuste dinámico de spread por volatilidad y profundidad, optimización para mantenerse Top 1, plug-and-play (cliente conecta su API y listo), y BYOK.'

'Modelo: USD 99 a 1,500/mes según volumen, más performance fee 5 a 10% para los whales.'

'Proyección: 25 merchants mid + 3 whales en 18 meses = 15 a 25 mil de MRR limpio. Cero competencia ML hoy en Colombia.'

'Por qué es la punta de lanza: los merchants se hablan entre ellos. Si uno gana 1% más al mes, los otros 50 quieren el bot.'""")

# ── Slide 14 — Proyección 24m ───────────────────────────────────────────────
def slide_14():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Proyección financiera 24 meses",
                    "Escenario conservador · 4 líneas en paralelo · BYOK = 0 costo variable")
    add_image(s, ASSETS / "charts/14_mrr_projection.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 14)
    set_notes(s, """[60s — Proyección financiera]

'Estos números son CONSERVADORES — quiero ser claro. Inversor sofisticado descuenta proyecciones por 50%, así que prefiero subprometer.'

Recorrer la curva:
- 'M3-M6: solo P2P generando — 5 a 6k MRR.'
- 'M12: P2P consolidado, hedging arrancando, remesas y freelancers entrando — 25k MRR.'
- 'M24: las 4 líneas corriendo, target 100 a 150k MRR — eso es 1.2 a 1.8 millones ARR.'

'Margen 70 a 85% típico SaaS. BYOK garantiza cero costo variable. La banda lima muestra el target M24.'

Si te preguntan: '¿qué pasa si una línea no escala?' → 'Por eso tenemos cuatro. Si quitas la peor — supongamos remesas, ciclo de venta más largo — sigues en 75 a 125k. La diversificación interna nos protege.'

Si te preguntan ARR: '1.2 a 1.8M ARR a 24 meses con burn modesto. Margen contribución 75% típico.'""")

# ── Slide 15 — Dónde estamos hoy ───────────────────────────────────────────
def slide_15():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Dónde estamos hoy",
                    "Cierre del núcleo · 4 productos en construcción paralela")

    items = [
        ("Pipeline L0–L7", "27 DAGs Airflow productivos, retraining semanal", LIME, "✓"),
        ("Dashboard Next.js", "8 páginas, 47 API routes, datos reales", LIME, "✓"),
        ("Smart v2.0 backtest", "+25.63% / Sharpe 3.35 / p = 0.006 (validado)", LIME, "✓"),
        ("Conexión Binance + MEXC", "operativa para ejecución (paper / testnet / live)", LIME, "✓"),
        ("Regime Gate live", "13/14 semanas mean-reverting bloqueadas Q1 2026", LIME, "✓"),
        ("Bot P2P MVP", "5 merchants beta — mes 1", AMBER, "→"),
        ("Hedging MVP", "diseño + 1 PYME piloto — mes 2-3", AMBER, "→"),
        ("Newsletter premium", "USD 99/mes — mes 1", AMBER, "→"),
    ]
    cols = 2
    for i, (title, desc, color, icon) in enumerate(items):
        col = i % cols
        row = i // cols
        cx = 0.5 + col * 6.45
        cy = 1.85 + row * 1.25
        add_rect(s, Inches(cx), Inches(cy), Inches(6.35), Inches(1.05),
                 fill=BG_PANEL, line=color, line_w=1.2)
        # icon
        ic = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.25), Inches(cy + 0.3),
                                Inches(0.45), Inches(0.45))
        ic.fill.solid(); ic.fill.fore_color.rgb = color
        ic.line.fill.background()
        add_text(s, Inches(cx + 0.25), Inches(cy + 0.3), Inches(0.45), Inches(0.45),
                 icon, size=14, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(cx + 0.85), Inches(cy + 0.15), Inches(5.4), Inches(0.45),
                 title, size=14, bold=True, color=WHITE)
        add_text(s, Inches(cx + 0.85), Inches(cy + 0.55), Inches(5.4), Inches(0.45),
                 desc, size=10, color=GRAY)

    add_footer(s, 15)
    set_notes(s, """[40s — Estado actual]

'Status check honesto: lo verde está hecho. Lo ámbar está en construcción.'

'Hecho: pipeline ML completo en producción, dashboard funcional, backtest auditado, exchanges conectados, regime gate operando en vivo.'

'En progreso: bot P2P MVP, hedging MVP, newsletter premium — los tres del primer trimestre del roadmap.'

'No vendemos vaporware. Vendemos producto + roadmap claro de comercialización.'""")

# ── Slide 16 — Roadmap 12 meses ─────────────────────────────────────────────
def slide_16():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Roadmap 12 meses",
                    "Q2 2026 → Q1 2027 · 4 swimlanes paralelos · ejecución secuencial inteligente")
    add_image(s, ASSETS / "diagrams/16_roadmap.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 16)
    set_notes(s, """[45s — Roadmap]

'Lectura por swimlane:'

'P2P arranca primero — beta en Q2, 25 clientes pagos en Q3, 50 + perf fees en Q4, expansión LATAM en Q1 2027. Cash flow rápido valida la tesis.'

'Hedging entra en Q3 con MVP, 10 PYMES piloto en Q4, primer banco partner en Q1 2027.'

'Remesas piloto en Q4 con 2 casas de cambio — ciclo de venta largo, por eso lo posicionamos detrás del cash flow del P2P.'

'App freelancers: lanzamiento beta cerrada en Q1 2027.'

Hitos en círculos lima abajo: 'M3 cierre v2.0 + 6k MRR, M9 25k MRR + Microsoft Founders Hub, M12 banco partner, M15 100-150k MRR + Series Seed.'

'Esto NO es un roadmap fantasioso — cada hito tiene OKR concreto y deadline.'""")

# ── Slide 17 — Hitos 90 días ────────────────────────────────────────────────
def slide_17():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Próximos 90 días — hitos críticos",
                    "Lo que vamos a entregar antes del 31 de Julio 2026")
    add_image(s, ASSETS / "diagrams/17_milestones_90.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 17)
    set_notes(s, """[40s — 90 días]

'Esto es lo que ustedes pueden medir. Si en 90 días no entregamos esto, fallamos.'

Mes 1 (Mayo): 'Cerrar Smart v2.0, firmar 5 merchants beta de P2P, deck para Microsoft listo, newsletter live.'

Mes 2 (Junio): '10 clientes pagos de P2P, primer piloto hedging PYME, aplicación a Microsoft for Startups, outreach a 5 casas de cambio.'

Mes 3 (Julio): 'MRR 5k+, créditos Azure asegurados, primer prospecto formal de casa de cambio.'

'Si en 90 días tienen un follow-up con nosotros, ya van a poder revisar contra esta lista. Esa es nuestra propuesta de accountability.'""")

# ── Slide 18 — ¿Qué necesitamos de Microsoft? ────────────────────────────────
def slide_18():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "¿Qué necesitamos de Microsoft?",
                    "4 pedidos concretos · ROI visible para ambos lados")

    items = [
        ("01", "Créditos Azure", "Founders Hub: hasta USD 150k para correr ML/OpenAI/Container Apps", BLUE),
        ("02", "Co-selling LATAM", "Introducción a clientes corporativos: bancos, casas de cambio, PYMES", LIME),
        ("03", "Soporte técnico", "Azure ML + Azure OpenAI · arquitecto solutions dedicado", AMBER),
        ("04", "Visibilidad LATAM", "Caso de éxito conjunto en eventos, blog, casos Microsoft for Startups", PINK),
    ]
    yt = 2.0
    for n, title, desc, color in items:
        # number circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(yt + 0.15),
                                  Inches(0.85), Inches(0.85))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        add_text(s, Inches(0.7), Inches(yt + 0.15), Inches(0.85), Inches(0.85),
                 n, size=20, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # body
        add_rect(s, Inches(1.85), Inches(yt), Inches(11), Inches(1.1),
                 fill=BG_PANEL, line=color, line_w=1.2)
        add_text(s, Inches(2.1), Inches(yt + 0.18), Inches(10.5), Inches(0.5),
                 title, size=18, bold=True, color=WHITE)
        add_text(s, Inches(2.1), Inches(yt + 0.6), Inches(10.5), Inches(0.5),
                 desc, size=12, color=GRAY)
        yt += 1.25

    add_footer(s, 18)
    set_notes(s, """[40s — Pedido a Microsoft]

'4 pedidos concretos, sin enredos:'

1. 'Créditos Azure: aplicación a Founders Hub, hasta 150 mil dólares. Esto cubre nuestra infra de ML por 18 meses.'

2. 'Co-selling LATAM: introducción a 3 clientes en bancos, casas de cambio, o PYMES grandes. Una sola intro nos puede convertir un piloto en contrato anual.'

3. 'Soporte técnico: arquitecto Azure ML + Azure OpenAI dedicado. Migración rápida y sin tropiezos.'

4. 'Visibilidad: caso de éxito documentado para Microsoft for Startups en LATAM — todos ganamos.'

'Es un partnership, no un cheque. Por eso vamos a la siguiente slide.'""")

# ── Slide 19 — ¿Qué ofrecemos a cambio? ──────────────────────────────────────
def slide_19():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "¿Qué ofrecemos a Microsoft?",
                    "Caso de éxito LATAM en fintech ML · pipeline de leads · cobranding")

    items = [
        ("Caso de éxito documentado", "Microsoft for Startups LATAM en fintech ML — story público con métricas reales", GREEN),
        ("Migración stack a Azure", "Container Apps + Azure ML + Azure OpenAI + Functions — arquitectura referencia", BLUE),
        ("Cobranding en marketing", "GlobalMinds powered by Microsoft Azure en website, decks, eventos LATAM", AMBER),
        ("Pipeline de leads", "Cada cliente que cerramos en bancos / casas de cambio / PYMES = lead corporativo Microsoft", PINK),
    ]
    yt = 2.0
    for title, desc, color in items:
        # check icon
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(yt + 0.2),
                                  Inches(0.7), Inches(0.7))
        circ.fill.solid(); circ.fill.fore_color.rgb = color
        circ.line.fill.background()
        add_text(s, Inches(0.7), Inches(yt + 0.2), Inches(0.7), Inches(0.7),
                 "✓", size=22, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_rect(s, Inches(1.7), Inches(yt), Inches(11.15), Inches(1.1),
                 fill=BG_PANEL, line=color, line_w=1.2)
        add_text(s, Inches(2), Inches(yt + 0.18), Inches(10.7), Inches(0.5),
                 title, size=18, bold=True, color=WHITE)
        add_text(s, Inches(2), Inches(yt + 0.6), Inches(10.7), Inches(0.5),
                 desc, size=12, color=GRAY)
        yt += 1.25

    add_footer(s, 19)
    set_notes(s, """[35s — Lo que ofrecemos]

'Esto no es asimétrico — Microsoft también gana:'

1. 'Caso de éxito en LATAM: ustedes están buscando casos emblemáticos para Microsoft for Startups en la región. Nosotros los aspiramos a ser ese caso.'

2. 'Migramos stack completo a Azure: Container Apps, Azure ML, Azure OpenAI, Functions. Arquitectura que pueden usar como referencia.'

3. 'Cobranding: GlobalMinds powered by Microsoft Azure visible en cada deck, evento, web.'

4. 'Pipeline de leads — cada banco o casa de cambio que cerramos como cliente, automáticamente entra en el pipeline corporativo de Microsoft. Win-win.'""")

# ── Slide 20 — Equipo (sin fotos) ───────────────────────────────────────────
def slide_20():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Equipo fundador",
                    "Stack técnico + tracción comercial · 18 meses construyendo")

    members = [
        {
            "name": "Pedro Sánchez Briceño",
            "role": "Co-founder & CTO",
            "tagline": "Senior Data Analyst / MLOps Engineer · Azure AI-102 certified",
            "creds": [
                "Aera Energy — Senior Data Analyst",
                "CRC — Análisis cuantitativo regulatorio",
                "CODALTEC — FACSAT-2 (proyecto satelital)",
                "Vault Insurance — MLOps en producción",
            ],
            "color": LIME,
        },
        {
            "name": "Freddy",
            "role": "Co-founder & CEO",
            "tagline": "Estrategia comercial · Red de contactos corporativos LATAM",
            "creds": [
                "Liderazgo comercial sector financiero",
                "Red activa en bancos y casas de cambio",
                "Ciclo completo de venta B2B y partnerships",
                "Comunicación ejecutiva con C-level",
            ],
            "color": BLUE,
        },
    ]
    for i, m in enumerate(members):
        cx = 0.5 + i * 6.45
        cy = 1.9
        add_rect(s, Inches(cx), Inches(cy), Inches(6.35), Inches(4.9),
                 fill=BG_PANEL, line=m["color"], line_w=1.6)
        # avatar circle with initials
        initials = "".join([w[0] for w in m["name"].split()[:2]]).upper()
        ic = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.4), Inches(cy + 0.4),
                                Inches(1.2), Inches(1.2))
        ic.fill.solid(); ic.fill.fore_color.rgb = m["color"]
        ic.line.fill.background()
        add_text(s, Inches(cx + 0.4), Inches(cy + 0.4), Inches(1.2), Inches(1.2),
                 initials, size=30, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        # name + role
        add_text(s, Inches(cx + 1.85), Inches(cy + 0.45), Inches(4.3), Inches(0.5),
                 m["name"], size=18, bold=True, color=WHITE)
        add_text(s, Inches(cx + 1.85), Inches(cy + 0.95), Inches(4.3), Inches(0.5),
                 m["role"], size=13, color=m["color"], italic=True)
        add_text(s, Inches(cx + 0.4), Inches(cy + 1.85), Inches(5.5), Inches(0.5),
                 m["tagline"], size=11, color=GRAY)

        # divider
        add_lime_underline(s, Inches(cx + 0.4), Inches(cy + 2.45), Inches(0.6))
        # creds
        cy0 = cy + 2.7
        for cred in m["creds"]:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx + 0.4), Inches(cy0 + 0.07),
                                     Inches(0.12), Inches(0.12))
            dot.fill.solid(); dot.fill.fore_color.rgb = m["color"]
            dot.line.fill.background()
            add_text(s, Inches(cx + 0.6), Inches(cy0), Inches(5.3), Inches(0.4),
                     cred, size=11, color=WHITE)
            cy0 += 0.45

    add_footer(s, 20)
    set_notes(s, """[35s — Equipo]

'Dos co-founders, perfiles complementarios:'

Pedro: 'Soy yo. CTO. Senior Data Analyst y MLOps Engineer con experiencia en Aera Energy, CRC, CODALTEC en el proyecto satelital FACSAT-2, y Vault Insurance. Certificado Azure AI-102. El motor de IA que vieron es mío de punta a punta.'

Freddy: 'Mi co-fundador. CEO. Lleva las relaciones comerciales — su red en bancos y casas de cambio en Colombia es lo que abre las puertas que un perfil técnico no puede abrir solo. Ciclo completo de venta B2B.'

'Equipo de 2 hoy. Roadmap incluye 1 senior engineer en Q3 cuando bot P2P escale.'""")

# ── Slide 21 — Tracción ─────────────────────────────────────────────────────
def slide_21():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Tracción y validación",
                    "18 meses construyendo · producción real · 4 oportunidades validadas")
    add_image(s, ASSETS / "diagrams/21_traction.png",
              Inches(0.4), Inches(1.5), width=Inches(12.5))
    add_footer(s, 21)
    set_notes(s, """[35s — Tracción]

'Cierro con los datos duros que un inversor mira:'

'18 meses de desarrollo. Backtest 2025 OOS: 25.63% con p-value 0.006. Sharpe 3.35. Win rate 82.4% en 34 trades. 4 productos en roadmap activo.'

'Y debajo: backtest auditado con metodología walk-forward, sistema en producción con datos reales — 25 servicios Docker, comparables internacionales validan precio entre 175 millones de euros y 12 mil millones de dólares, regime gate bloqueó 13 de 14 semanas mean-reverting en Q1 2026 protegiéndonos del crash.'

'No estoy vendiéndoles humo — les estoy vendiendo evidencia.'""")

# ── Slide 22 — Call to action + disclaimer ─────────────────────────────────
def slide_22():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Próximos pasos",
                    "Lo que pedimos esta semana · cómo seguir conversando")

    # 3 CTA cards
    ctas = [
        ("Microsoft", "Aplicar a Microsoft for Startups Founders Hub · reunión técnica con Azure ML team",
         BLUE),
        ("Contactos comerciales", "Piloto 60 días GRATUITO · 3 clientes elegidos (1 PYME, 1 merchant P2P, 1 casa de cambio)",
         LIME),
        ("Inversores", "Acceso temprano a ronda Seed Q4 2026 · ticket mínimo USD 25k · valuación SAFE post-money TBD",
         AMBER),
    ]
    yt = 1.85
    for title, desc, color in ctas:
        add_rect(s, Inches(0.5), Inches(yt), Inches(12.35), Inches(1.0),
                 fill=BG_PANEL, line=color, line_w=1.5)
        add_rect(s, Inches(0.5), Inches(yt), Inches(2.5), Inches(1.0),
                 fill=color, line=None, line_w=0)
        add_text(s, Inches(0.5), Inches(yt + 0.25), Inches(2.5), Inches(0.5),
                 title, size=15, bold=True, color=BG_DARK, align=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)
        add_text(s, Inches(3.2), Inches(yt + 0.2), Inches(9.5), Inches(0.65),
                 desc, size=12, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
        yt += 1.15

    # contact
    add_rect(s, Inches(0.5), Inches(5.4), Inches(12.35), Inches(0.8),
             fill=BG_PANEL, line=LIME, line_w=1)
    add_text(s, Inches(0.7), Inches(5.55), Inches(12), Inches(0.5),
             "📧  pspedroelias96@gmail.com    ·    🔗  globalminds.ai (en construcción)    ·    📅  Agendemos follow-up esta semana",
             size=13, color=WHITE, bold=True)

    # disclaimer (prominent)
    add_rect(s, Inches(0.5), Inches(6.35), Inches(12.35), Inches(0.65),
             fill=BG_DARK, line=RED, line_w=1)
    add_text(s, Inches(0.7), Inches(6.4), Inches(12), Inches(0.55),
             "⚠  GlobalMinds NO constituye asesoría financiera. Resultados pasados NO garantizan resultados futuros. El uso de los productos implica riesgo de pérdida de capital.",
             size=10, color=RED, bold=True)

    add_footer(s, 22)
    set_notes(s, """[40s — Cierre]

'Tres asks claros — uno por audiencia:'

1. 'Microsoft: queremos aplicar a Founders Hub esta semana y agendar reunión técnica con el equipo de Azure ML.'

2. 'A los contactos comerciales: ofrecemos un piloto de 60 días GRATUITO a 3 clientes elegidos por ustedes — 1 PYME, 1 merchant P2P, 1 casa de cambio.'

3. 'A inversores: acceso temprano a ronda Seed planificada para Q4 2026. Ticket mínimo 25k. Si están interesados, conversación esta semana.'

'Mi correo está aquí. Mi compromiso: respondo en menos de 24 horas a cualquier follow-up de esta sala.'

Disclaimer (importante leer si la audiencia es estricta): 'Por compliance: GlobalMinds no constituye asesoría financiera. Performance pasada no garantiza futura. Hay riesgo de pérdida de capital. Esto es prospecto de partnership, no recomendación de inversión.'

Pausa de 3 segundos. Mirar a la audiencia. Sonreír. 'Gracias.'""")

# ── Slide 23 — Sources & References ─────────────────────────────────────────
def slide_23():
    s = prs.slides.add_slide(blank)
    set_bg(s)
    add_title_block(s, "Apéndice — Fuentes y referencias",
                    "Cifras públicas y verificables citadas en este deck")

    # 2 columns
    col1 = [
        ("Banco de la República — Remesas",
         "USD 13,098 M en 2025 (récord) · banrep.gov.co/es/estadisticas/remesas"),
        ("Banco de la República — Borradores 1058",
         "17% de empresas no financieras con cobertura FX · banrep.gov.co/es/borrador-1058"),
        ("Visa adquiere Kantox",
         "€175M (2025) · visa.com/newsroom · TechCrunch · El País Negocios"),
        ("Bre-B — Pagos Instantáneos Colombia",
         ">100M de llaves activas · banrep.gov.co/bre-b"),
        ("Backtest GlobalMinds 2025",
         "Smart Simple v2.0 OOS · Walk-forward · summary_2025.json (auditado)"),
    ]
    col2 = [
        ("Wise — Capitalización pública",
         "USD 11B · wise.com/investor-relations · LSE: WISE"),
        ("Remitly — Capitalización pública",
         "USD 3B · NASDAQ: RELY"),
        ("Deel — Última ronda",
         "USD 12B · pitchbook.com · TechCrunch"),
        ("Binance P2P — Spread y fees",
         "0% fees · binance.com/p2p · MetaReporte (calculadora)"),
        ("XTX Markets",
         "Líder global ML market making FX · xtxmarkets.com"),
    ]

    for col_idx, items in enumerate([col1, col2]):
        cx = 0.5 + col_idx * 6.45
        yt = 1.85
        for title, desc in items:
            add_rect(s, Inches(cx), Inches(yt), Inches(6.35), Inches(0.95),
                     fill=BG_PANEL, line=GRAY_DIM, line_w=0.5)
            add_text(s, Inches(cx + 0.3), Inches(yt + 0.13), Inches(6), Inches(0.45),
                     title, size=12, bold=True, color=LIME)
            add_text(s, Inches(cx + 0.3), Inches(yt + 0.5), Inches(6), Inches(0.45),
                     desc, size=10, color=WHITE)
            yt += 1.05

    # bottom note
    add_text(s, Inches(0.5), Inches(7.0), Inches(12), Inches(0.3),
             "Lista completa con URLs activas en SOURCES.md (incluido en el paquete del deck).",
             size=10, color=GRAY, italic=True)

    add_footer(s, 23)
    set_notes(s, """[Apéndice — backup para preguntas]

Esta slide NO se presenta — está para responder preguntas del tipo 'de dónde sacaste ese dato'.

Tener a mano el SOURCES.md por si piden URLs específicas.

Datos sensibles a defender:
- 17% cobertura FX: estudio Banrep, fecha del paper. Reconocer que es lo más reciente público.
- USD 13B remesas: cifra Banrep 2025. Récord histórico.
- Kantox/Visa €175M: confirmado por Visa Newsroom + TechCrunch + Diario El País 2025.
- Comparables Wise/Remitly/Deel: información pública de mercados bursátiles + Pitchbook.""")


# ─── Build all ────────────────────────────────────────────────────────────────
slide_01(); slide_02(); slide_03(); slide_04(); slide_05()
slide_06(); slide_07(); slide_08(); slide_09(); slide_10()
slide_11(); slide_12(); slide_13(); slide_14(); slide_15()
slide_16(); slide_17(); slide_18(); slide_19(); slide_20()
slide_21(); slide_22(); slide_23()

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"✓ {OUT} ({OUT.stat().st_size // 1024} KB)")
print(f"  {len(prs.slides)} slides")
